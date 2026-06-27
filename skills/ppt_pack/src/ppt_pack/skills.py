"""PPT 生成 skill 实现 — 四阶段流水线.

4 个 skill:
  - data_collect: 收集数据（stub：返回模板数据）
  - outline:     写大纲（stub：返回固定 JSON 结构）
  - layout:      排版（stub：返回 Markdown）
  - export:      导出（stub：返回文件路径字符串）

跟 agentvet_pack / crawler_pack 相同套路:
  每个 skill = 1 个 class, 实现 core/skill.py 的 Skill ABC.
"""
from __future__ import annotations

import logging
from typing import Any

from core.skill import Skill, SkillManifest

_log = logging.getLogger(__name__)


# ── Skill 1: data_collect ────────────────────────────────────────────────

class DataCollectSkill(Skill):
    """收集数据 — stub 返回模板数据, 后续可接入真实数据源.

    输入: {"topic": "主题名", "sources": [...]}  — sources 可选
    输出: {"ok": True, "data": {"title": ..., "sections": [...], "stats": [...]}}
    """

    def __init__(self) -> None:
        super().__init__(SkillManifest(
            name="data_collect",
            api_version="1.0",
            description="收集 PPT 所需数据（stub 模板, 后续接真实源）",
            tags=("ppt", "collect"),
        ))

    def run(self, input_data: dict) -> dict:
        topic = input_data.get("topic", "未指定主题")
        output_file = input_data.get("output_file", "")
        return {
            "ok": True,
            "topic": topic,
            "output_file": output_file,
            "data": {
                "title": topic,
                "sections": ["背景", "现状分析", "解决方案", "总结展望"],
                "stats": [
                    {"label": "覆盖率", "value": "85%"},
                    {"label": "增长率", "value": "+12%"},
                ],
                "sources": ["内部数据", "行业报告"],
            },
        }


# ── Skill 2: outline ─────────────────────────────────────────────────────

class OutlineSkill(Skill):
    """写大纲 — stub 返回固定 JSON 结构, 供下游 Layout 消费.

    输入: {"data": {"title": ..., "sections": [...]}}  或 {"topic": "..."}
    输出: {"ok": True, "outline": {"title": ..., "slides": [...]}}
    """

    def __init__(self) -> None:
        super().__init__(SkillManifest(
            name="outline",
            api_version="1.0",
            description="生成 PPT 大纲（stub 固定结构）",
            tags=("ppt", "outline"),
        ))

    def run(self, input_data: dict) -> dict:
        data = input_data.get("data", {})
        topic = data.get("title") or input_data.get("topic", "未指定主题")
        sections = data.get("sections", ["概述", "详情", "结论"])
        return {
            "ok": True,
            "outline": {
                "title": topic,
                "slides": [
                    {
                        "index": i + 1,
                        "title": sec,
                        "bullet_points": [f"{sec} 要点一", f"{sec} 要点二"],
                    }
                    for i, sec in enumerate(sections)
                ],
            },
        }


# ── Skill 3: layout ──────────────────────────────────────────────────────

class LayoutSkill(Skill):
    """排版 — stub 返回 Markdown, 不依赖 python-pptx.

    输入: {"outline": {"title": ..., "slides": [...]}}
    输出: {"ok": True, "markdown": "...", "slide_count": N}
    """

    def __init__(self) -> None:
        super().__init__(SkillManifest(
            name="layout",
            api_version="1.0",
            description="PPT 排版渲染（stub 返回 Markdown）",
            tags=("ppt", "layout"),
        ))

    def run(self, input_data: dict) -> dict:
        outline = input_data.get("outline", {})
        title = outline.get("title", "未命名")
        slides = outline.get("slides", [])
        md_lines = [f"# {title}", ""]
        for s in slides:
            md_lines.append(f"## {s['title']}")
            for bp in s.get("bullet_points", []):
                md_lines.append(f"- {bp}")
            md_lines.append("")
        return {"ok": True, "markdown": "\n".join(md_lines), "slide_count": len(slides)}


# ── Skill 4: export ──────────────────────────────────────────────────────

class ExportSkill(Skill):
    """导出 — 用 python-pptx 生成真实 .pptx 文件.

    输入: {"markdown": "..."} 或 {"layouts": [...]}（兼容两种上游）
          "output": "out.pptx" — 可选输出路径
    输出: {"ok": True, "file": "out.pptx", "size_kb": N, "slide_count": N}
    """

    def __init__(self) -> None:
        super().__init__(SkillManifest(
            name="export",
            api_version="1.0",
            description="导出 PPT 文件（python-pptx 真实生成）",
            tags=("ppt", "export"),
        ))

    def run(self, input_data: dict) -> dict:
        markdown = input_data.get("markdown", "")
        layouts = input_data.get("layouts", [])

        if layouts:
            slides = layouts
        elif markdown:
            slides = self._parse_markdown(markdown)
        else:
            return {"ok": False, "error": "need markdown or layouts"}

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import os

            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 宽屏
            prs.slide_height = Inches(7.5)

            for s in slides:
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                if title and s.get("title"):
                    title.text = s["title"]
                if s.get("content"):
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.text = s["content"][:500]

            output_file = input_data.get("output", "output.pptx")
            prs.save(output_file)
            file_size = os.path.getsize(output_file)
            return {
                "ok": True,
                "file": output_file,
                "size_kb": file_size // 1024,
                "slide_count": len(prs.slides),
            }
        except Exception as exc:
            return {"ok": False, "error": str(exc)}

    def _parse_markdown(self, md: str) -> list[dict]:
        """简易 markdown → slides.

        # 标题 → 封面 slide
        ## 标题 → 内容 slide（标题行）
        其他行 → 追加到当前 slide 内容区
        """
        slides: list[dict] = []
        for line in md.split("\n"):
            if line.startswith("# ") and not line.startswith("## "):
                slides.append({"title": line[2:].strip(), "content": ""})
            elif line.startswith("## ") and slides:
                slides[-1]["content"] += line[3:].strip() + "\n"
            elif slides:
                slides[-1]["content"] += line.strip() + "\n"
        return slides


# 注册入口 (给 Pool 用)
def register_all(pool) -> int:
    """把 4 个 skill 全部注册到 pool. 返回注册数."""
    n = 0
    for cls in (DataCollectSkill, OutlineSkill, LayoutSkill, ExportSkill):
        try:
            pool.register(cls())
            n += 1
        except Exception as exc:  # noqa: BLE001
            _log.warning("register %s failed: %s", cls.__name__, exc)
    return n
